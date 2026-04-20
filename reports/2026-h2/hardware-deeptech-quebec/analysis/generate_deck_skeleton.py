"""Generate the Day Zero deck skeleton as a .pptx file.

This script is reproducible — rerun it whenever the slide list in
SLIDES needs to evolve. The skeleton is meant as a navigable starting
point for the team, not a final design. Visual design and brand polish
happen at the end of the project.

Requirements
------------
python-pptx (not in the default project pyproject; install in a local
venv):

    python3 -m venv /tmp/deckvenv
    /tmp/deckvenv/bin/pip install python-pptx
    /tmp/deckvenv/bin/python reports/2026-h2/hardware-deeptech-quebec/analysis/generate_deck_skeleton.py

The generator writes to:
    reports/2026-h2/hardware-deeptech-quebec/drafts/day-zero-deck-skeleton.pptx

Design
------
16:9 widescreen. Navy + gray + accent orange. Cover, section dividers,
analytical slides (2x2 grid with labeled boxes), and content slides.
Speaker notes carry the full Day Zero context for each analytical slide
so that the deck stays coherent even when opened in isolation.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# __file__ = .../ecosystem/reports/2026-h2/hardware-deeptech-quebec/analysis/generate_deck_skeleton.py
# parents[0]=analysis  [1]=hardware-...  [2]=2026-h2  [3]=reports  [4]=ecosystem
REPO_ROOT = Path(__file__).resolve().parents[4]
OUTPUT = (
    REPO_ROOT
    / "reports"
    / "2026-h2"
    / "hardware-deeptech-quebec"
    / "drafts"
    / "day-zero-deck-skeleton.pptx"
)

# ---------------------------------------------------------------------
# Palette and sizing
# ---------------------------------------------------------------------

NAVY = RGBColor(0x0A, 0x25, 0x40)
NAVY_SOFT = RGBColor(0x1C, 0x3D, 0x5A)
GRAY_DARK = RGBColor(0x37, 0x41, 0x51)
GRAY_MID = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ---------------------------------------------------------------------
# Primitive helpers
# ---------------------------------------------------------------------


def add_text_box(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font_size=12,
    font_color=GRAY_DARK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_multiline_text_box(
    slide, x, y, w, h, paragraphs, *, font_size=11, font_color=GRAY_DARK,
):
    """paragraphs: list of (text, is_bold) or just strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, para in enumerate(paragraphs):
        if isinstance(para, tuple):
            text, bold = para
        else:
            text, bold = para, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return tb


def add_filled_rect(slide, x, y, w, h, fill_color, *, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_footer(slide, slide_num: int, total: int, *, section_label: str | None = None):
    left = (
        f"Day Zero draft · v0.2 · Hardware & Deep Tech au Québec · brouillon interne"
        if not section_label
        else f"Day Zero · {section_label}"
    )
    add_text_box(
        slide,
        Inches(0.3),
        Inches(7.15),
        Inches(10),
        Inches(0.25),
        left,
        font_size=8,
        font_color=GRAY_MID,
    )
    add_text_box(
        slide,
        Inches(12.0),
        Inches(7.15),
        Inches(1.1),
        Inches(0.25),
        f"{slide_num} / {total}",
        font_size=8,
        font_color=GRAY_MID,
        align=PP_ALIGN.RIGHT,
    )


def set_speaker_notes(slide, notes: str):
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes


# ---------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------


def build_cover(prs, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_filled_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(1.2), NAVY)
    add_filled_rect(slide, Inches(0), Inches(6.3), Inches(SLIDE_W_IN), Inches(1.2), GRAY_LIGHT)
    add_text_box(
        slide,
        Inches(1),
        Inches(2.2),
        Inches(11.33),
        Inches(1.2),
        "Hardware et Deep Tech au Québec",
        font_size=44,
        font_color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(3.5),
        Inches(11.33),
        Inches(0.8),
        "État des lieux, dynamiques et leviers pour la prochaine décennie",
        font_size=20,
        font_color=GRAY_DARK,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.8),
        Inches(11.33),
        Inches(0.45),
        "Un rapport conjoint de Quebec Tech + Réseau Capital",
        font_size=16,
        font_color=NAVY_SOFT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(5.3),
        Inches(11.33),
        Inches(0.4),
        "(Conseil de l'Innovation du Québec — co-signataire en discussion)",
        font_size=11,
        font_color=GRAY_MID,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(6.55),
        Inches(11.33),
        Inches(0.35),
        "Day Zero deck · v0.2 · 2026-04-17 · brouillon interne · à ne pas diffuser",
        font_size=10,
        font_color=GRAY_MID,
        align=PP_ALIGN.CENTER,
    )


def build_section_divider(prs, num: int, title: str, subtitle: str, slide_num: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN), NAVY)
    add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11.33),
        Inches(1.2),
        f"Section {num}",
        font_size=22,
        font_color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(3.3),
        Inches(11.33),
        Inches(1.2),
        title,
        font_size=52,
        font_color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.6),
        Inches(11.33),
        Inches(1),
        subtitle,
        font_size=16,
        font_color=GRAY_LIGHT,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, slide_num, total, section_label=f"Section {num} · {title}")


def build_analytical_slide(
    prs,
    title: str,
    question: str,
    donnees: str,
    visuel: str,
    takeaway: str,
    caveats: str,
    status: str,
    slide_num: int,
    total: int,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    add_filled_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(0.9), NAVY)
    add_text_box(
        slide,
        Inches(0.4),
        Inches(0.18),
        Inches(12.5),
        Inches(0.6),
        title,
        font_size=22,
        font_color=WHITE,
        bold=True,
    )

    # 2x2 grid of boxes. Layout area: y=1.1 to y=6.9 (5.8 tall), x=0.3 to 13.0 (12.7 wide)
    # Gap = 0.2 in
    box_w = Inches(6.15)
    box_h = Inches(2.75)
    top_y = Inches(1.15)
    bot_y = Inches(4.05)
    left_x = Inches(0.35)
    right_x = Inches(6.75)

    def draw_box(x, y, label, body):
        add_filled_rect(slide, x, y, box_w, box_h, GRAY_LIGHT)
        # Label
        add_text_box(
            slide,
            x,
            y,
            box_w,
            Inches(0.35),
            f"  {label.upper()}",
            font_size=9,
            font_color=ORANGE,
            bold=True,
        )
        # Body
        tb = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.4), box_w - Inches(0.2), box_h - Inches(0.5)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = body
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_DARK

    draw_box(left_x, top_y, "Question analytique", question)
    draw_box(right_x, top_y, "Données / sources", donnees)
    draw_box(left_x, bot_y, "Visuel (placeholder)", visuel)
    draw_box(right_x, bot_y, "Takeaway attendu · caveats", f"{takeaway}\n\nCaveats: {caveats}")

    # Status strip at very bottom
    add_text_box(
        slide,
        Inches(0.35),
        Inches(6.85),
        Inches(12.5),
        Inches(0.25),
        f"STATUS: {status}",
        font_size=9,
        font_color=NAVY_SOFT,
        bold=True,
    )

    add_footer(slide, slide_num, total)

    # Speaker notes carry the raw Day Zero context so the deck stays
    # coherent when opened standalone.
    notes = (
        f"SLIDE: {title}\n\n"
        f"QUESTION ANALYTIQUE:\n{question}\n\n"
        f"DONNÉES / SOURCES:\n{donnees}\n\n"
        f"VISUEL (PLACEHOLDER):\n{visuel}\n\n"
        f"TAKEAWAY ATTENDU:\n{takeaway}\n\n"
        f"CAVEATS:\n{caveats}\n\n"
        f"STATUS: {status}"
    )
    set_speaker_notes(slide, notes)


def build_content_slide(
    prs, title: str, bullets: list[str], slide_num: int, total: int, *, speaker_notes: str = ""
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(0.9), NAVY)
    add_text_box(
        slide,
        Inches(0.4),
        Inches(0.18),
        Inches(12.5),
        Inches(0.6),
        title,
        font_size=22,
        font_color=WHITE,
        bold=True,
    )

    # Bullets area
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY_DARK

    add_footer(slide, slide_num, total)
    if speaker_notes:
        set_speaker_notes(slide, speaker_notes)


# ---------------------------------------------------------------------
# Slide definitions
# ---------------------------------------------------------------------

# Each entry is a dict with a 'type' field: 'cover', 'section', 'content',
# or 'analytical'. Additional fields depend on the type.

SLIDES = [
    {"type": "cover"},
    # ---------- Section 1 — Cadrage ----------
    {
        "type": "section",
        "num": 1,
        "title": "Cadrage",
        "subtitle": "TL;DR, contexte, définitions, méthodologie",
    },
    {
        "type": "content",
        "title": "Ce qu'il faut retenir",
        "bullets": [
            "Environ [X] compagnies deep-tech au Québec — lasers, puces, robots, dispositifs médicaux, équipement spatial. Ancrées par INO, EXFO, MDA, Kinova, Nord Quantique.",
            "Un réveil depuis 2024, mais modeste. ~18 nouvelles compagnies/trimestre en 2025–début 2026 vs ~12 pré-COVID. Quantique et matériel adjacent bougent; photonique traditionnelle reste plate.",
            "Le vent GenAI n'a pas porté le matériel. Création software × 2.2 depuis 2022; création hardware plate. Deux trajectoires opposées.",
            "Le vrai problème n'est pas de créer, c'est de faire grandir. Capital patient manque, demande domestique diffuse, opérateurs expérimentés rares.",
            "[Recommandation phare à formuler après analyse] — leviers plausibles : capital dédié, procurement stratégique, immigration accélérée, renforcement des anchors.",
        ],
        "notes": "Slide à réécrire en toute fin de projet — les 5 messages doivent sortir de l'analyse, pas la précéder. Ton grand public : chaque message compréhensible sans contexte technique.",
    },
    {
        "type": "analytical",
        "title": "Pourquoi ce rapport, pourquoi maintenant",
        "question": "Quelle est la fenêtre de pertinence de ce rapport ?",
        "donnees": "Contexte global : CHIPS Act (US 2022), Chips Act (EU 2023), réponse canadienne limitée. Contexte QC : héritage photonique INO/EXFO, pôle quantique émergent, contraste post-GenAI SW/HW. Contexte politique : cycle électoral QC fin 2026.",
        "visuel": "Timeline horizontale 2017–2026 avec événements-clés (COVID, CHIPS Act, vague GenAI, élections QC) superposée à une courbe de formation hard-tech.",
        "takeaway": "Fenêtre de pertinence claire : les politiques deep-tech se durcissent à l'international; le Québec a les ancres mais pas le momentum; les élections ouvrent l'espace discursif.",
        "caveats": "Ne pas surcharger le contexte politique électoral — rapport grand public, pas partisan.",
        "status": "Narratif prêt; viz à construire.",
    },
    {
        "type": "analytical",
        "title": "Deep tech, hard tech : de quoi parle-t-on ?",
        "question": "Qu'est-ce que ce rapport appelle 'deep tech' et 'hard tech' ? Comment distingue-t-on ces compagnies des autres start-ups technologiques ?",
        "donnees": "Deep tech = bâti sur avancée scientifique/ingénierie substantielle, cycles R-D longs, capital-intense. Hard tech = sous-ensemble à produit physique au cœur. Ce rapport couvre deep tech large (9 piliers).",
        "visuel": "Schéma en cercles concentriques (deep tech large → hard tech → 9 piliers), avec liste courte 'INCLUS / EXCLU' et 3 exemples publics QC par colonne.",
        "takeaway": "Le lecteur grand public doit, en 30 secondes, comprendre où est la frontière entre ce qu'on compte et ce qu'on exclut (SaaS pur, IA pure, fintech).",
        "caveats": "Définitions stabilisées dans analysis/taxonomy-and-definitions.md. Renvoi explicite pour le lecteur qui veut creuser.",
        "status": "Définitions arrêtées; viz à construire.",
    },
    {
        "type": "analytical",
        "title": "Les 9 piliers du deep tech québécois",
        "question": "Comment le deep tech québécois se décompose-t-il concrètement ?",
        "donnees": "Table 9 lignes : photonique, quantique, semi, robotique, spatial, matériaux avancés, medtech/biotech-hw, cleantech-hw, agri-tech-hw. Définition en une phrase + un exemple public par pilier (EXFO, Nord Quantique, IBM Bromont, Kinova, MDA, NanoXplore, Medtronic, Lion Electric, Semios).",
        "visuel": "Grille 3×3 de tuiles avec icône, nom du pilier, phrase-définition, logo/nom d'une compagnie publique québécoise.",
        "takeaway": "Le lecteur comprend en 30 secondes la diversité des piliers et peut associer chaque pilier à un nom concret — démystifie la notion abstraite de 'deep tech'.",
        "caveats": "Cohérence critique avec analysis/taxonomy-and-definitions.md. Tout changement de nomenclature ici doit être répercuté là.",
        "status": "Contenu prêt; viz à construire.",
    },
    {
        "type": "analytical",
        "title": "Comment on a fait (méthodologie en 30 secondes)",
        "question": "Comment les chiffres ont-ils été obtenus ? Quelle confiance peut-on y accorder ?",
        "donnees": "Point de départ : REQ. Classification hybride mots-clés + codes CAE. Croisement Dealroom (QT) et PitchBook (RC). Trois niveaux de confiance : HAUTE / MOYENNE / BASSE.",
        "visuel": "Infographie horizontale en 4 étapes (REQ → classification hybride → croisement sources → fiabilité tiers) avec icônes.",
        "takeaway": "On est transparent sur les plages de confiance. Les compagnies incorporées au fédéral (CBCA) sont un angle mort à combler via Dealroom/PitchBook.",
        "caveats": "73 % des classifications historiques reposaient sur CAE seul. Importance d'annoncer les plages de confiance.",
        "status": "Texte prêt. Accord RC critique pour la partie PitchBook — à obtenir dans les 2 premières semaines.",
    },
    # ---------- Section 2 — Portrait ----------
    {
        "type": "section",
        "num": 2,
        "title": "Portrait",
        "subtitle": "Combien, quoi, où, depuis quand, quelle taille",
    },
    {
        "type": "analytical",
        "title": "Combien de compagnies deep-tech au Québec",
        "question": "Quel est l'univers de compagnies deep-tech au Québec en 2026, et comment est-il borné ?",
        "donnees": "SILVER.REQ_PRODUCT_CLASSIFICATION filtré DEEPTECH + cross-ref GOLD.STARTUP_REGISTRY. Extension de l'analyse existante INTERNAL-hardware-photonics-req-2026 au DEEPTECH complet.",
        "visuel": "Grand chiffre avec fourchette, sous-titré par décomposition par pilier.",
        "takeaway": "Fournir le chiffre-référence que tout le monde citera, avec une fourchette honnête plutôt qu'un nombre précis. Décomposer plutôt que donner un chiffre unique.",
        "caveats": "Compagnies CBCA sous-détectées. Compagnies sans description ou site web sous-détectées.",
        "status": "Analyse à étendre — existant couvre photonique+matériel adjacent; inclure quantique, robotique, spatial, matériaux, biotech-hw.",
    },
    {
        "type": "analytical",
        "title": "Répartition par pilier",
        "question": "Comment se répartit cet univers entre les 9 piliers deep-tech ?",
        "donnees": "SILVER.REQ_PRODUCT_CLASSIFICATION.MATCHED_SIGNALS — compter par token. Validation manuelle sur échantillon par bucket.",
        "visuel": "Barres horizontales triées décroissant, ou treemap respectant la hiérarchie mère-fille de la taxonomie.",
        "takeaway": "La photonique et le matériel adjacent dominent historiquement; quantique et robotique sont petits en volume mais croissent; spatial sous-représenté malgré anchors médiatisés.",
        "caveats": "Un NEQ peut matcher plusieurs tokens (biotech + medtech typique) — distinguer DISTINCT NEQ de total occurrences.",
        "status": "Partielle — Q5 couvre post-2024; rétrospectif long à produire.",
    },
    {
        "type": "analytical",
        "title": "Concentration géographique",
        "question": "Où sont concentrées ces compagnies au Québec ?",
        "donnees": "SILVER.REQ_PRODUCT_CLASSIFICATION.HQ_CITY (extrait regex d'adresse REQ). Agrégation par RMR via taxonomy/geographies.yaml.",
        "visuel": "Carte choroplèthe du Québec colorée par densité + overlay des institutions-phares (INO Québec, universités MTL/SHE).",
        "takeaway": "Grand Montréal domine en volume; Grand Québec a probablement la plus forte densité relative en photonique grâce à l'INO et Laval.",
        "caveats": "Extraction HQ_CITY par regex = bruitée; 'unknown' non-négligeable.",
        "status": "Données prêtes; viz à construire.",
    },
    {
        "type": "analytical",
        "title": "Pyramide d'âge",
        "question": "La population deep-tech québécoise est-elle neuve, mature, ou mixte ?",
        "donnees": "INCORPORATION_YEAR avec buckets 2026, 2020–2025, 2010–2019, 2000–2009, <2000.",
        "visuel": "Histogramme par année d'incorporation avec ligne moyenne/médiane; superposition software/AI pour contraste si pertinent.",
        "takeaway": "Pyramide probablement plus aplatie que software — moins de création annuelle mais plus de maturité historique.",
        "caveats": "INCORPORATION_YEAR = année légale, pas nécessairement début d'activité commerciale. Pas d'info cessation d'activité.",
        "status": "Données prêtes; analyse simple à lancer.",
    },
    {
        "type": "analytical",
        "title": "Distribution par taille d'employés",
        "question": "Quelle est la taille typique d'une compagnie deep-tech québécoise ?",
        "donnees": "N_EMPLOYES et EMP_MIN numérique. Buckets REQ : 1–5, 6–10, 11–20, 21–50, 51–100, 101–250, 251–500, 501+.",
        "visuel": "Pyramide ou bar chart des tranches d'employés, couleur par pilier si lisible.",
        "takeaway": "Queue longue à droite — beaucoup de petites (1–20), peu au-dessus de 250. Scaling gap visible dès les effectifs.",
        "caveats": "Tranches autodéclarées à l'inscription, non mises à jour en temps réel. 'Non déclaré' = bucket significatif.",
        "status": "Données prêtes.",
    },
    {
        "type": "analytical",
        "title": "Modèle d'affaires (B2B / B2C / B2G)",
        "question": "Quelle est la distribution B2B / B2C / B2G chez les deep-tech québécoises ?",
        "donnees": "shared_ecosystem.qt_schema.companies.b2b_b2c (via Dealroom) + compléments manuels top 50 via sites web et LinkedIn.",
        "visuel": "Donut 3 segments ou 100 % stacked bar par pilier.",
        "takeaway": "Hard tech est massivement B2B ou B2G. Très peu de B2C pur. Contraste avec software consumer; cycles de vente longs.",
        "caveats": "Couverture Dealroom variable sur deep-tech. Beaucoup de small-scale n'ont pas de champ b2b_b2c renseigné.",
        "status": "Partielles; couverture Dealroom à quantifier.",
    },
    {
        "type": "analytical",
        "title": "Phares établis (anchors)",
        "question": "Quelles compagnies établies servent de 'bordure' à l'écosystème et autour desquelles gravite la formation neuve ?",
        "donnees": "Listes publiques seulement : EXFO, Coherent/II-VI ex-MPB (2021), MDA, CAE, Lion Electric, Nord Quantique, Kitco, AddÉnergie, GHGSat. Annonces médiatisées, sites web.",
        "visuel": "Tableau ou logos mosaïque organisé par pilier, avec année de fondation, HQ, statut (privé / public / acquis).",
        "takeaway": "Documenter la 'gravité' institutionnelle — le Québec a des anchors crédibles dans plusieurs verticales, atout pour spinoffs.",
        "caveats": "Inclure uniquement ce qui est vérifiable publiquement. 15–25 noms max.",
        "status": "Liste à construire — validation QT sur qui inclure.",
    },
    # ---------- Section 3 — Dynamiques ----------
    {
        "type": "section",
        "num": 3,
        "title": "Dynamiques",
        "subtitle": "Formation, cohortes récentes, financement, sorties, trajectoires",
    },
    {
        "type": "analytical",
        "title": "Taux de formation trimestriel 2017–2026",
        "question": "Comment a évolué le rythme d'incorporation de compagnies deep-tech au Québec sur une décennie ?",
        "donnees": "SILVER.REQ_PRODUCT_CLASSIFICATION filtré DEEPTECH, groupé par DATE_TRUNC('quarter', DATE_IMMATRICULATION). L'analyse existante donne : pré-COVID 12.1/q → COVID 15.1/q → post-CHIPS 11.3/q → 2024 12.2/q → 2025–2026 18.2/q.",
        "visuel": "Ligne trimestrielle avec 4 zones ombrées (pré-COVID, COVID, post-CHIPS, récent) et annotations d'événements-clés.",
        "takeaway": "Soubresaut récent (~50 % au-dessus de la baseline pré-COVID) mais petits volumes absolus. Confidence: Low (MSG-HARDTECH-01).",
        "caveats": "Échantillons trimestriels à un chiffre pour sous-secteurs. 73 % via CAE seulement.",
        "status": "Viz existante en draft; à élargir au DEEPTECH complet.",
    },
    {
        "type": "analytical",
        "title": "Le soubresaut 2025 — signal ou bruit ?",
        "question": "Le pic apparent 2025–début 2026 est-il statistiquement robuste ou un artéfact de petit échantillon ?",
        "donnees": "Même source que slide précédente. Test : retirer 3–4 incorporations du trimestre pic et voir si la tendance tient. Cross-ref MEDTEQ+, INO si accès.",
        "visuel": "Zoom sur 2024–2026 avec bandes de confiance (bootstrap rough), annotation 'en l'absence de X incorporations, la tendance disparaît'.",
        "takeaway": "Honnêteté : à monitorer, pas à publier comme fait. Signal probable mais requiert 3–4 trimestres additionnels.",
        "caveats": "Reprendre mot à mot la mise en garde MSG-HARDTECH-01.",
        "status": "Analyse à faire; calcul de sensibilité simple.",
    },
    {
        "type": "analytical",
        "title": "Photonique pure vs matériel adjacent",
        "question": "La photonique québécoise, avec sa base institutionnelle (INO), génère-t-elle des nouvelles compagnies, ou la croissance vient-elle d'ailleurs ?",
        "donnees": "Analyse existante : photonique pure flat ~2.0–2.2/q 2017–2026; croissance dans optique/imagerie (7/q en 2025) et matériel adjacent (9/q en 2025).",
        "visuel": "3 lignes superposées : photonique core, optique/imagerie, matériel adjacent. Temps sur X, formation/trimestre sur Y.",
        "takeaway": "La photonique pure est plate malgré l'INO, EXFO, etc. La base institutionnelle ne génère PAS la vague de spinoffs qu'on attendrait — à creuser.",
        "caveats": "MSG-HARDTECH-02 (Confidence: Medium). Classification sous-catégorie repose sur MATCHED_SIGNALS.",
        "status": "Analyse existante; à élargir DEEPTECH complet.",
    },
    {
        "type": "analytical",
        "title": "La cohorte post-2024",
        "question": "Qui sont les nouveaux entrants hard-tech post-2024, et sont-ils prometteurs (HIGH) ou bruit (LOW CAE-only) ?",
        "donnees": "Q5 diagnostic (pipelines/validation/diagnostics/Q5_req_post2024_hardtech.sql) déjà écrit. Sections 4–5 : décomposition par sous-catégorie et tier de confiance.",
        "visuel": "Tableau des nouveaux entrants par pilier × confiance (HIGH / MEDIUM / LOW), top 10 HIGH listés nommément si publiquement visibles.",
        "takeaway": "La cohorte post-2024 penche vers HIGH-confidence quantique et robotique. Donner un chiffre 'propre' et un chiffre 'brut' pour montrer la marge d'incertitude.",
        "caveats": "Q5 dépend de SILVER.REQ_PRODUCT_CLASSIFICATION rafraîchi dans Snowflake. Top 50 HIGH à valider manuellement avant citation.",
        "status": "SQL prêt; rouler dans Snowsight et trier CSV.",
    },
    {
        "type": "analytical",
        "title": "Contraste software vs hard-tech post-GenAI",
        "question": "Le boom GenAI 2023–2025 a-t-il entraîné le hard-tech, ou creusé l'écart ?",
        "donnees": "Comparaison taux de formation software-AI vs deep-tech. Analyse existante INTERNAL-genai-impact-req-2026 : création tech-product × 2.2 depuis 2022 (software) vs hardware plat.",
        "visuel": "2 lignes indexées à 100 en 2017 : software-AI et hard-tech. Évolution claire de l'écart.",
        "takeaway": "Le vent GenAI n'a pas porté le hard-tech. Divergence structurelle qui justifie des instruments de politique publique DÉDIÉS. Message-pivot du rapport.",
        "caveats": "Classification software-AI a ses propres bruits.",
        "status": "Analyse existante partielle; à formuler clairement la comparaison.",
    },
    {
        "type": "analytical",
        "title": "Financement hard-tech au Québec",
        "question": "Quel volume et quels montants de capital vont au hard-tech québécois, vs software ?",
        "donnees": "PitchBook/CVCA via RC — AGRÉGATS SEULEMENT (gouvernance). Dealroom funding QC deep-tech. Annonces publiques cross-référencées.",
        "visuel": "Deux histogrammes empilés côte-à-côte : 'Software-AI' vs 'Hard-tech' pour deals et montants. Timeline 2020–2025.",
        "takeaway": "Hard-tech reçoit fraction disproportionnellement petite du capital levé, deals plus rares mais plus gros en moyenne.",
        "caveats": "Source licenciée → agrégats seulement. Accord RC BLOQUANT.",
        "status": "BLOQUANT sur accord RC. Rédiger demande formelle.",
    },
    {
        "type": "analytical",
        "title": "Sorties notables",
        "question": "Quelles acquisitions et IPOs ont marqué le hard-tech québécois sur la dernière décennie ?",
        "donnees": "Annonces publiques (Bloomberg, communiqués, presse spécialisée) — sources publiques uniquement. PitchBook exits pour vérification si accord RC.",
        "visuel": "Timeline horizontale avec 8–15 sorties marquantes, taille des bulles = valeur si publique.",
        "takeaway": "Quelques sorties significatives dominées par acquisitions par stratégiques US. Peu d'IPOs. Sortie canadienne/québécoise rare → implication écosystème risque.",
        "caveats": "Sorties privées non annoncées hors tableau. Valeurs souvent non divulguées.",
        "status": "Liste à bâtir depuis sources publiques.",
    },
    {
        "type": "analytical",
        "title": "Trajectoires par pilier",
        "question": "Quels piliers deep-tech montent, stagnent, ou reculent ?",
        "donnees": "Time-series de formation par pilier (9). Optionnel : croisement avec données de financement.",
        "visuel": "Small multiples — 9 mini-lignes en grille 3×3, normalisées, avec badge '↑ / → / ↓'.",
        "takeaway": "Carte synthétique des dynamiques — quantique et robotique montent, photonique plate, semi quasi-inexistant, spatial stable. Oriente sections 6 et 7.",
        "caveats": "Petits échantillons par pilier — traiter comme qualitatif plutôt que précis.",
        "status": "Analyse à faire.",
    },
    # ---------- Section 4 — Acteurs ----------
    {
        "type": "section",
        "num": 4,
        "title": "Acteurs",
        "subtitle": "Institutions, capital, accélérateurs, universités, talent",
    },
    {
        "type": "analytical",
        "title": "Institutions de recherche et de transfert",
        "question": "Quelle est la base institutionnelle sur laquelle s'appuie le hard-tech québécois ?",
        "donnees": "Publiques : INO, MEDTEQ+, Écotech Québec, IVADO, Mila, Institut Quantique (Sherbrooke), PINQ² (Bromont), Prompt, CRIQ, CNRC Boucherville. Listes de membres si accord.",
        "visuel": "Carte géographique QC avec institutions positionnées, taille proportionnelle à budget/effectifs publics.",
        "takeaway": "Base institutionnelle dense et géographiquement distribuée (Québec-MTL-Sherbrooke-Bromont). Peu d'écosystèmes canadiens ont autant d'infrastructures.",
        "caveats": "Budget/effectifs publics quand divulgués; sinon qualitatif.",
        "status": "Texte prêt; viz carte à construire.",
    },
    {
        "type": "analytical",
        "title": "Financement public et non-dilutif",
        "question": "Quels instruments publics soutiennent le hard-tech ? Lesquels sont lisibles par des opérateurs ?",
        "donnees": "Publiques : SR&ED, PARI/IRAP, programmes MEIE, MITACS, FRQNT, Investissement Québec, Canada SIF, Stratégie quantique, CanExport. MEIE rapport.",
        "visuel": "Tableau matriciel 'instrument × stade × pilier applicable' avec codage couleur lisibilité/pertinence.",
        "takeaway": "Beaucoup d'instruments, mais lisibilité et coordination déficientes. Pas d'équivalent CHIPS Act — lacune structurelle.",
        "caveats": "Informations datent vite (budgets fédéraux changent) — dater la slide.",
        "status": "Inventaire à construire.",
    },
    {
        "type": "analytical",
        "title": "Investisseurs privés actifs en hard-tech",
        "question": "Quels fonds VC / PE investissent dans le hard-tech québécois ? Lacunes à certains stades ?",
        "donnees": "PitchBook/CVCA via RC — agrégats. Sites des fonds : Inovia, Real Ventures, Amplitude, BDC Capital, Investissement Québec, Fonds FTQ, Fondaction, CDPQ, Cycle Capital, Amorchem.",
        "visuel": "Bar chart top 15 investisseurs par # deals hard-tech QC 2020–2025; annotation du stade.",
        "takeaway": "Capital seed disponible, capital scale-up (Série B+) rare et souvent étranger. Goulot classique canadien.",
        "caveats": "Licensed → agrégats seulement. 'Actif' = au moins N deals — seuil à justifier.",
        "status": "BLOQUANT sur accord RC.",
    },
    {
        "type": "analytical",
        "title": "Accélérateurs et incubateurs à orientation hard-tech",
        "question": "Quels programmes québécois ont une spécialisation ou une ouverture réelle au hard-tech ?",
        "donnees": "Listes publiques : Centech (Polytechnique), District 3 (Concordia), TandemLaunch, Cycle Momentum, MEDTEQ+, Prompt. Cohortes récentes (sites web).",
        "visuel": "Tableau : nom / affiliation / cohorte/an / verticale / stade / exemple de diplômé hard-tech.",
        "takeaway": "Quelques programmes solides (Centech, TandemLaunch, MEDTEQ+), mais écosystème accélération plus tourné vers software.",
        "caveats": "Cohortes varient; données 2024–2025 partielles.",
        "status": "Inventaire à construire.",
    },
    {
        "type": "analytical",
        "title": "Universités et spinoffs",
        "question": "Les universités QC génèrent-elles des spinoffs hard-tech à un rythme compétitif ?",
        "donnees": "Bureaux de transfert (Univalor, SOVAR, Aligo) — rapports annuels. Listes de spinoffs universitaires (Sherbrooke publique). Cross-ref REQ via nom/adresse/fondateur.",
        "visuel": "Barres horizontales par université : spinoffs hard-tech 2015–2025, avec ratio sur budget recherche.",
        "takeaway": "Sherbrooke punches above its weight en quantique. McGill forte en photonique/semi. Polytechnique + UdeM forts en robotique/aéro. Conversion recherche → compagnies sous-optimale.",
        "caveats": "Données spinoffs non-normalisées entre universités. Federal CBCA spinoffs sous-détectés.",
        "status": "Données à collecter; lourd mais important.",
    },
    {
        "type": "analytical",
        "title": "Talent : ingénieurs, PhDs, opérateurs",
        "question": "Le bassin de talent hard-tech québécois est-il suffisant pour soutenir une accélération ?",
        "donnees": "StatCan ESCM/EAMT — diplômés STEM. Tableau MEIE sur talent. Rapports GSER/StartupBlink sur rang Québec.",
        "visuel": "3 indicateurs : (1) diplômés STEM/an, (2) PhDs/an, (3) 'operators avec expérience scale-up hard-tech' (qualitatif).",
        "takeaway": "Pool de diplômés solide. Opérateurs expérimentés (C-level hard-tech ayant vécu scale-up) rares — importation nécessaire, enjeu immigration.",
        "caveats": "Qualitatif pour la 3e métrique — reconnaître la limite.",
        "status": "Données publiques prêtes; qualitatif à argumenter.",
    },
    # ---------- Section 5 — Benchmarks ----------
    {
        "type": "section",
        "num": 5,
        "title": "Benchmarks",
        "subtitle": "8 régions par archétype : petits champions, politique industrielle, grandes économies, émergent, domestique",
    },
    {
        "type": "analytical",
        "title": "Panorama des 8 benchmarks",
        "question": "Comment le Québec se positionne-t-il globalement face aux 8 régions de comparaison ?",
        "donnees": "GSER 2024–2025, StartupBlink Global 2025, OCDE R-D. Rapports nationaux : Vinnova (SE), Israel Innovation Authority, A*STAR (SG), KISTEP (KR), BMWK (DE), NSF (US), NCBR (PL), ISED+StatCan (CA).",
        "visuel": "Grand tableau/heatmap 8 régions + QC × 6 indicateurs-clés — PIB/habitant, R-D privée %PIB, compagnies deep-tech/M hab, capital deep-tech par habitant, universités top-200, index politique industrielle.",
        "takeaway": "Le QC a les ressources institutionnelles de rivaliser avec petits champions, mais volume de formation et capital sous le seuil critique. Position intermédiaire.",
        "caveats": "Définitions 'deep tech' varient par source. Recul 1–2 ans sur données nationales.",
        "status": "Collecte documentaire à lancer.",
    },
    {
        "type": "analytical",
        "title": "Petits champions deep-tech : Suède, Israël, Singapour",
        "question": "Que font ces 3 petites économies que le QC n'a pas encore fait, et qu'ont-elles en commun ?",
        "donnees": "SE : Vinnova, RISE, Chalmers Ventures, ancrages Ericsson/AstraZeneca. IL : IIA, Technion spinoffs, Yozma model, pipeline militaire→civil (8200). SG : A*STAR, NRF, SGInnovate, Temasek, ciblage top-down.",
        "visuel": "3 colonnes parallèles — un 'portrait' par pays (compagnies, capital, instruments-phares, leçon). En bas : 'Ce que les 3 partagent.'",
        "takeaway": "3 ingrédients communs : (1) anchors industriels forts, (2) capital public de CROISSANCE (pas seulement seed), (3) demande domestique captive (défense, santé, fonds souverain).",
        "caveats": "Contextes politiques/sécuritaires très différents (Israël). Pas tout transposable.",
        "status": "Recherche documentaire à faire.",
    },
    {
        "type": "analytical",
        "title": "Politique industrielle ciblée : Corée du Sud",
        "question": "Comment la Corée est-elle devenue une puissance semi/matériaux en une génération, et qu'est-ce qui est transposable au QC ?",
        "donnees": "KISTEP, K-Chips Act (2023), Samsung/SK Hynix disclosures, KDB reports. Littérature sur modèle chaebols et évolution deep-tech.",
        "visuel": "Timeline 1980–2025 des dépenses R-D publiques + événements-clés; encart 'chiffres-clés' comparés au Canada.",
        "takeaway": "Modèle coréen : coordination top-down massive + champions nationaux + éducation technique ciblée. Transposable à échelle réduite : ciblage sous-sectoriel (quantique, photonique), partenariats structurants, financement long-horizon.",
        "caveats": "Échelle et contexte culturel très différents. Composante 'champion national' controversée démocratiquement.",
        "status": "Recherche documentaire à faire.",
    },
    {
        "type": "analytical",
        "title": "Grandes économies : Allemagne et États-Unis",
        "question": "Que retenir des deux plus gros modèles (US CHIPS Act, DE Mittelstand + SPRIND/DTEC) et qu'est-ce qui survit à l'adaptation pour le QC ?",
        "donnees": "US : White House CHIPS reports, NSF, NIST, fabs annoncées, IRA cleantech. DE : BMWK, SPRIND, DTEC.Bw, Mittelstand deep-tech (Carl Zeiss, Trumpf).",
        "visuel": "2 colonnes — US ('big bet subventions') et DE ('dense réseau mittelstand'). Tableau comparatif + 'Ce que le QC peut importer à son échelle.'",
        "takeaway": "Ni modèle US (52 G$ CHIPS) ni Mittelstand ne se reproduisent à l'identique, mais 2 leçons : (1) ciblage d'infrastructures manquantes (US : fab), (2) financement institutionnel long-terme du réseau existant (DE : Fraunhofer). CRIQ et INO = Fraunhofer canadiens potentiels.",
        "caveats": "CHIPS Act 3 ans de recul seulement. Mittelstand produit d'un siècle.",
        "status": "Recherche documentaire à faire.",
    },
    {
        "type": "analytical",
        "title": "Émergent : Pologne (cas de rattrapage)",
        "question": "Une économie européenne de rattrapage post-intégration EU peut-elle inspirer une trajectoire d'accélération du QC ?",
        "donnees": "NCBR reports, Polish Development Fund (PFR) deep-tech initiatives. OCDE/Eurostat R-D.",
        "visuel": "Trajectoire 2004–2025 (R-D, compagnies deep-tech) avec annotations inflexions politiques; courbe parallèle QC pour comparaison.",
        "takeaway": "La Pologne montre qu'une décennie de politique industrielle cohérente peut multiplier la base deep-tech par 3–5×. Le QC est plus mature mais a un déficit d'accélération similaire. Horizon = décennie, pas mandat.",
        "caveats": "Fonds EU structurants non-disponibles au Canada. Contexte post-soviétique spécifique.",
        "status": "Recherche documentaire à faire — angle inédit dans rapports canadiens existants.",
    },
    {
        "type": "analytical",
        "title": "Domestique : Québec vs reste du Canada",
        "question": "Quelle est la place du Québec dans le deep tech canadien, et où souffre-t-il par rapport à Ontario et CB ?",
        "donnees": "StatCan R-D par province. ISED KSBS. Ontario Business Registry (accès à confirmer). GSER/StartupBlink rankings. Rapports MaRS, Communitech (ON), BC Tech.",
        "visuel": "Carte canadienne choroplèthe + tableau comparatif QC / ON / CB / AB / Atlantique sur 6 indicateurs.",
        "takeaway": "Ontario plus gros en volume absolu (pop 1.8× QC) avec cluster semi Waterloo + biotech Toronto. QC surreprésenté en photonique, quantique, cleantech, aérospatial par habitant. Le QC = province hardware-centric du Canada.",
        "caveats": "Définitions provinciales asymétriques. Comparaison 'toutes choses égales' difficile.",
        "status": "Recherche à faire; Ontario Business Registry à explorer.",
    },
    {
        "type": "analytical",
        "title": "CHIPS Act, Chips Act EU, Stratégie quantique CA : la pause canadienne",
        "question": "Quels effets observés des grandes politiques industrielles récentes, et qu'est-ce qui manque côté canadien ?",
        "donnees": "White House CHIPS reports + CRS. European Commission dashboards. Conseil consultatif canadien sur les semi (2023). Stratégie quantique CA (360M$).",
        "visuel": "Timeline 2020–2026 avec montants engagés et résultats à 3 ans; colonne 'Canada' visiblement plus pâle.",
        "takeaway": "Le Canada n'a pas de CHIPS Act. Instruments existants utiles mais sous-dimensionnés et dispersés. Le QC a les leviers pour agir unilatéralement sur photonique et quantique.",
        "caveats": "CHIPS Act 3 ans de recul — appréciation préliminaire.",
        "status": "Recherche documentaire à faire.",
    },
    {
        "type": "analytical",
        "title": "Capital par stade : QC vs pairs",
        "question": "À quels stades de financement le Québec souffre-t-il le plus en deep-tech vs pairs directement comparables ?",
        "donnees": "PitchBook/CVCA via RC. Stratification Seed/A/B/Growth pour deep-tech QC vs ON, US moyen, IL, SE.",
        "visuel": "4 barres par région (Seed, A, B, Growth) — QC vs 3–4 régions comparables.",
        "takeaway": "Seed correct, Série A correct, Séries B et Growth critiquement manquants pour deep-tech au QC. Goulot canadien amplifié par intensité capitalistique.",
        "caveats": "Stages Dealroom vs PitchBook peuvent différer — normaliser via taxonomy/stages.yaml. Licensed → agrégats.",
        "status": "BLOQUANT sur accord RC.",
    },
    # ---------- Section 6 — Défis ----------
    {
        "type": "section",
        "num": 6,
        "title": "Défis et lacunes",
        "subtitle": "Scaling, capex, talent, supply chain, fragmentation",
    },
    {
        "type": "analytical",
        "title": "Le gap de scaling",
        "question": "Combien de compagnies hard-tech québécoises passent du stade seed au scaleup (10M$+ ARR / 50+ employés) ?",
        "donnees": "Compter dans GOLD.STARTUP_REGISTRY + données financement les compagnies franchissant les seuils. Comparer ratio formation → scaleup au software.",
        "visuel": "Funnel — formation / premier tour / Série A / scaleup — pour hard-tech vs software.",
        "takeaway": "Funnel hard-tech beaucoup plus étroit au sommet. Capital patient et demande domestique manquent.",
        "caveats": "Cohortes petites → valeurs bruitées.",
        "status": "Analyse à faire.",
    },
    {
        "type": "analytical",
        "title": "Intensité capitalistique et capex",
        "question": "Le capital disponible au QC est-il adapté à l'intensité capex du hard-tech (équipements, fab, certifications) ?",
        "donnees": "Estimations capex par pilier (littérature publique, comparables). Taille moyenne Série A hard-tech vs software au QC.",
        "visuel": "Barres 'capex requis pour break-even' par pilier vs 'taille moyenne Série A au QC'.",
        "takeaway": "Écart structurel — les rounds typiques QC ne suffisent pas pour MVP hardware dans la moitié des piliers. Rationale pour instruments publics ou programmes MatchCap.",
        "caveats": "Estimations capex varient énormément intra-pilier.",
        "status": "Recherche + analyse à faire.",
    },
    {
        "type": "analytical",
        "title": "Pénurie de talents-clés",
        "question": "Quels profils manquent le plus, et comment ça bloque la croissance ?",
        "donnees": "Enquêtes MEIE/Emploi-Québec sur pénurie STEM. Témoignages (entretiens qualitatifs à conduire).",
        "visuel": "Matrix 2×2 — profil × sévérité de pénurie (qualitatif).",
        "takeaway": "VP Engineering avec scaling hard-tech, ingénieurs process semi, PhDs quantique appliqué, operators EU/US avec réseaux — les plus rares.",
        "caveats": "Qualitatif; à valider par entretiens.",
        "status": "Entretiens à planifier.",
    },
    {
        "type": "analytical",
        "title": "Capacité manufacturière et supply chain",
        "question": "Le QC peut-il faire scale ses compagnies hard-tech localement, ou dépend-il de sous-traitance externe (US, Asie) ?",
        "donnees": "Inventaire des foundries/ateliers propres (INO fab, CMC Microsystems, AAC). Rapports stratégie quantique CA + semi.",
        "visuel": "Carte des capacités manufacturières au Canada, heatmap par pilier.",
        "takeaway": "Foundry quantique (Bromont — IBM/PINQ²) remarquable. Globalement dépendance forte hors-QC pour fab semi, assembly, certification. Vulnérabilité géopolitique.",
        "caveats": "Ne pas surestimer 'faiblesse' — certains choix de dépendance sont rationnels (scale).",
        "status": "Inventaire à construire.",
    },
    {
        "type": "analytical",
        "title": "Fragmentation fédéral-provincial",
        "question": "Les instruments fédéraux et provinciaux s'articulent-ils en ensemble cohérent, ou créent-ils confusion ?",
        "donnees": "Cartographie des programmes (slide précédente) remixée par niveau de gouvernement. Entretiens qualitatifs avec opérateurs.",
        "visuel": "Diagramme en flux de la trajectoire d'un opérateur cherchant financement/support — nombre de portes d'entrée, temps moyen.",
        "takeaway": "Fragmentation crée friction — fondateur typique consulte 4–6 programmes. Opportunité pour guichet unique QC.",
        "caveats": "Qualitatif à l'appui d'entretiens.",
        "status": "Entretiens à planifier.",
    },
    # ---------- Section 7 — Recommandations ----------
    {
        "type": "section",
        "num": 7,
        "title": "Recommandations",
        "subtitle": "Placeholders — à formuler après analyse",
    },
    {
        "type": "content",
        "title": "Principes directeurs",
        "bullets": [
            "Ciblage sous-sectoriel, pas soutien générique 'tech'.",
            "Capital patient + demande domestique, les deux simultanés.",
            "Miser sur les anchors existants (INO, MEDTEQ+, IQ) plutôt que recréer.",
            "Coordonner fédéral-provincial, pas rivaliser.",
        ],
        "notes": "Principes à valider/reformuler après analyse. Nés des patterns observés dans les benchmarks et les défis identifiés.",
    },
    {
        "type": "content",
        "title": "Recommandation 1 (placeholder) — Capital matériel",
        "bullets": [
            "[PLACEHOLDER] Créer un fonds spécifique hard-tech (matching provincial + Investissement Québec + BDC).",
            "[PLACEHOLDER] Ticket-size alignés sur intensité capex du secteur (Séries B 25–75 M$).",
            "[PLACEHOLDER] Horizon de liquidité 10–15 ans.",
            "[PLACEHOLDER] Sous-secteurs prioritaires à déterminer par analyse.",
        ],
        "notes": "Formulation finale dépend du résultat de la section 6 (Capex vs Série A actuelle).",
    },
    {
        "type": "content",
        "title": "Recommandation 2 (placeholder) — Talent / immigration",
        "bullets": [
            "[PLACEHOLDER] Programme d'immigration accéléré pour opérateurs hard-tech expérimentés.",
            "[PLACEHOLDER] Partenariats avec anchors (INO, MEDTEQ+) pour ciblage.",
            "[PLACEHOLDER] Alignement avec stratégie provinciale en immigration économique.",
        ],
        "notes": "Formulation finale dépend du résultat de la section 6 (pénurie talents).",
    },
    {
        "type": "content",
        "title": "Recommandation 3 (placeholder) — Institutions-phares",
        "bullets": [
            "[PLACEHOLDER] Financement accru des bureaux de transfert universitaires QC avec mandat spinoffs hard-tech.",
            "[PLACEHOLDER] Fonds de proof-of-concept pré-incorporation.",
            "[PLACEHOLDER] Alignement avec INO, MEDTEQ+, Écotech Québec.",
        ],
        "notes": "Formulation finale dépend du résultat de la section 4 (acteurs) et 6 (lacunes).",
    },
    {
        "type": "content",
        "title": "Recommandation 4 (placeholder) — Demande domestique / procurement",
        "bullets": [
            "[PLACEHOLDER] Programmes de procurement stratégique (Hydro-Québec, Santé publique, transport).",
            "[PLACEHOLDER] Cadre d'approvisionnement auprès de fournisseurs hard-tech locaux.",
            "[PLACEHOLDER] Inspiré modèles Israël / Singapour.",
        ],
        "notes": "Formulation finale dépend des benchmarks petits champions (SE, IL, SG).",
    },
    {
        "type": "content",
        "title": "Roadmap et suivi",
        "bullets": [
            "[PLACEHOLDER] Priorisation 3-3-3 : 3 actions sur 12 mois, 3 sur 36 mois, 3 sur 5 ans.",
            "[PLACEHOLDER] Indicateurs de suivi quantitatifs (compagnies deep-tech, capital levé, exits, talent).",
            "[PLACEHOLDER] Gouvernance du suivi annuel.",
        ],
        "notes": "À remplir en fin de draft.",
    },
    # ---------- Section 8 — Annexe ----------
    {
        "type": "section",
        "num": 8,
        "title": "Annexe",
        "subtitle": "Sources, méthodologie détaillée, taxonomie, limites, reproductibilité, glossaire",
    },
    {
        "type": "content",
        "title": "Sources de données",
        "bullets": [
            "REQ (Registraire des entreprises du Québec) — public, extrait 2026-04.",
            "SILVER.REQ_PRODUCT_CLASSIFICATION — dérivé, construit par stage 31 classifier.",
            "GOLD.STARTUP_REGISTRY — registre consolidé QT + RC (Dealroom + Harmonic/PitchBook).",
            "Dealroom via QT — licencié, agrégats seulement.",
            "PitchBook via RC — licencié, agrégats seulement (pending accord formel).",
            "StatCan, ISED, MEIE — public.",
            "Rapports externes : GSER, StartupBlink, BDC, Deloitte, Stanford, MTB (voir insights/reports-external/).",
        ],
        "notes": "Slide à finaliser au moment du livrable. Mettre à jour les dates de coupure.",
    },
    {
        "type": "analytical",
        "title": "Méthodologie détaillée (1/3) — Classification hybride",
        "question": "Comment une compagnie-type voyage-t-elle à travers le classifier pour être étiquetée 'deep tech' ?",
        "donnees": "40+ patterns keywords sur DESC_LOWER et DESC_ALL. Codes CAE (boost ciblés 2851, 3361, 3674, 3827, 3740). Filtre service (regex exclusion) évite 3,757 faux positifs. Score → tier HIGH/MEDIUM/LOW.",
        "visuel": "Flowchart 'voyage d'une compagnie' — description + CAE → filtres → score → tier. Pédagogique.",
        "takeaway": "Classification reproductible et transparente. Chaque compagnie a un score et un tier, pas une étiquette binaire.",
        "caveats": "73% des classifications historiques reposaient sur CAE seul. Le filtre service élimine l'essentiel du bruit mais pas tout.",
        "status": "Référentiel prêt (pipelines/transforms/silver/31_req_product_classification.sql + analysis/taxonomy-and-definitions.md).",
    },
    {
        "type": "analytical",
        "title": "Méthodologie détaillée (2/3) — Définitions par pilier",
        "question": "Comment chaque pilier est-il défini de façon opérationnelle ?",
        "donnees": "Table 9 piliers × 5 colonnes — définition une phrase, 2-3 mots-clés principaux, codes CAE dominants, 2 exemples publics, caveats-clés.",
        "visuel": "Table dense 9 lignes × 5 colonnes; design clair avec codes CAE en petit.",
        "takeaway": "Le lecteur peut vérifier comment une compagnie qu'il connaît serait classifiée, et identifier si notre classification diverge de son intuition.",
        "caveats": "La taxonomie est v0.1 — révisions possibles via branche taxonomy/ du repo.",
        "status": "Contenu prêt (analysis/taxonomy-and-definitions.md section 2). Viz à construire.",
    },
    {
        "type": "analytical",
        "title": "Méthodologie détaillée (3/3) — Matching et gouvernance",
        "question": "Comment les sources multiples (REQ, Dealroom, PitchBook) sont-elles unifiées, et quelles règles de gouvernance appliquées aux données sensibles ?",
        "donnees": "Matching : bridge NEQ, fuzzy name ≥ 0.85, manual review whitelist. Gouvernance : agrégats seulement pour sources licenciées, seuil minimum 5 enregistrements.",
        "visuel": "Diagramme de flux sources → unification → publication, avec annotations gouvernance.",
        "takeaway": "La méthodologie de matching est documentée et reproductible. La gouvernance protège les données sensibles sans empêcher la production d'insights.",
        "caveats": "Compagnies incorporées fédéralement (CBCA) sous-détectées — angle mort structurel.",
        "status": "Référentiel prêt (pipelines/transforms/entity_resolution/ + DATA-GOVERNANCE.md).",
    },
    {
        "type": "analytical",
        "title": "Taxonomie employée",
        "question": "Quelle est la taxonomie de secteurs utilisée dans ce rapport, et comment se rattache-t-elle à la taxonomie canonique du repo ?",
        "donnees": "taxonomy/sectors.yaml (repo partagé) + taxonomy/startup-criteria.yaml. Périmètre du rapport = DEEPTECH + parties de HEALTHTECH et CLEANTECH + partie d'AGRITECH.",
        "visuel": "Arbre taxonomique 2 niveaux (parent → enfants) avec codes colorés 'inclus / partiellement inclus / exclu'.",
        "takeaway": "Le périmètre est clair et décidable : on peut pour chaque compagnie répondre oui/non à 'est-ce dans ce rapport ?'.",
        "caveats": "Taxonomie v0.1 — évolution prévue via branche taxonomy/ du repo.",
        "status": "Données prêtes (YAML). Viz à construire.",
    },
    {
        "type": "content",
        "title": "Limitations connues",
        "bullets": [
            "73% des classifications REQ historiques reposent sur CAE seul — prudence sur trajectoires fines de sous-secteurs.",
            "Compagnies incorporées fédéralement (CBCA) manquantes dans REQ — angle mort systématique.",
            "REQ capte l'incorporation, pas l'activité commerciale — des 'shells' peuvent être comptés.",
            "Small sample sizes par pilier → pas de tests statistiques fins applicables.",
            "Absence (pending) de données financement RC/PitchBook pour plusieurs slides — version dégradée si non obtenue.",
            "Classification de sous-catégorie repose sur tokens keyword — biais connus (compagnies sans description, anglicismes).",
            "Compagnies acquises/absorbées peuvent disparaître du suivi si réorganisation juridique.",
        ],
        "notes": "Section de transparence. À lire en miroir des takeaways de chaque slide — où chaque limitation s'applique-t-elle le plus ?",
    },
    {
        "type": "content",
        "title": "Reproductibilité",
        "bullets": [
            "Branche : report/hardware-deeptech-quebec",
            "Code : reports/2026-h2/hardware-deeptech-quebec/analysis/ + pipelines/",
            "Diagnostic post-2024 : pipelines/validation/diagnostics/Q5_req_post2024_hardtech.sql",
            "Classification silver : pipelines/transforms/silver/31_req_product_classification.sql",
            "Insights agrégés : insights/2026-h2/hardware-deeptech-*.md",
            "Taxonomie de référence : reports/2026-h2/hardware-deeptech-quebec/analysis/taxonomy-and-definitions.md",
        ],
        "notes": "À finaliser au moment du livrable. Objectif : un lecteur technique peut refaire tous les chiffres.",
    },
    {
        "type": "content",
        "title": "Équipe, remerciements, glossaire",
        "bullets": [
            "Équipe projet : lead QT (Étienne Bernard), lead RC (à nommer), lead CIQ (à nommer si in), reviewers, contributors.",
            "Remerciements : partenaires ayant partagé listes/données (INO, MEDTEQ+, Écotech Québec si confirmés).",
            "Glossaire : deep tech, hard tech, photonique, quantique, scale-up, CAE, NEQ, CBCA, Série A/B, SR&ED, IRAP, Mittelstand.",
            "Version : v0.2 · Data cutoff : 2026-04 · Contact : etienne@quebectech.com",
        ],
        "notes": "À finaliser en toute fin de projet.",
    },
]


# ---------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    total = len(SLIDES)
    slide_num = 0

    for defn in SLIDES:
        slide_num += 1
        t = defn["type"]
        if t == "cover":
            build_cover(prs, total)
        elif t == "section":
            build_section_divider(
                prs,
                defn["num"],
                defn["title"],
                defn.get("subtitle", ""),
                slide_num,
                total,
            )
        elif t == "content":
            build_content_slide(
                prs,
                defn["title"],
                defn["bullets"],
                slide_num,
                total,
                speaker_notes=defn.get("notes", ""),
            )
        elif t == "analytical":
            build_analytical_slide(
                prs,
                defn["title"],
                defn["question"],
                defn["donnees"],
                defn["visuel"],
                defn["takeaway"],
                defn["caveats"],
                defn["status"],
                slide_num,
                total,
            )
        else:
            raise ValueError(f"Unknown slide type: {t}")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {slide_num} slides to {OUTPUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
